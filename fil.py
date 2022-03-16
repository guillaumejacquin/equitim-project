import plotly.express as px
import pandas as pd
import numpy as np
from pathlib import Path

df = pd.DataFrame({
        "time": pd.date_range("1-Jul-2021", periods=200, freq="1H"),
        "desired_light": np.random.choice(["ON", "OFF"], 200),
    })

# build a figure
fig = px.bar(df, x="time", color="desired_light", height=400, title="Restaurant bills")
# export fingure as a static file: https://plotly.com/python/static-image-export/
fname = Path.cwd().joinpath("SO.png")
fig.write_image(fname)

# create a powerpoint from an image: https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
pic = slide.shapes.add_picture(str(fname), left, top)


prs.save(Path.cwd().joinpath('SO.pptx'))
