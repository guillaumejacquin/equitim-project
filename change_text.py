from cmd import Cmd
from dataclasses import dataclass
from typing import List
from unicodedata import name

from pyparsing import col
from calculs.abac import abac
from pptx import Presentation
from PIL import Image
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import os
from pptx.util import Pt

from copy import deepcopy
from pptx.table import Table, _Row, _Column, _Cell

def elementsToReplaceDegressivite(Class, shapes):
    #si degressif on supprime les balises
    if Class.BAC_is_degressif != "":
        Class.desonndr = ""
        Class.longuephrase = ""
        Class.SDBAC = ""
        Class.PDINSM = ""
        Class.ETPDI = ""

#TOUS LES REMPLACEMENTS 
def elementsToReplaceRemplacement(Class, shapes):
    #mise en avant car transforme en bac

    replace_text({'<NDRTES>': Class.desonndr}, shapes)
    replace_text({'<DESONNDR>': Class.desonndr}, shapes)

    replace_text({'<longuephrase>': Class.longuephrase}, shapes)
    replace_text({'<SDBAC>': Class.SDBAC}, shapes)
    replace_text({'<PDINSM>': Class.PDINSM}, shapes)
    replace_text({'<TEST>': Class.ETPDI}, shapes) 

    replace_text({'<EBAC>': Class.EBAC}, shapes) 

    replace_text({'<NOM>':  Class.Nom}, shapes)
    replace_text({'<droit>':  Class.Droit}, shapes)
    replace_text({'<ISIN>':  Class.Isin}, shapes)
    replace_text({'<émission>':  Class.Emission_affichage}, shapes)
    replace_text({'<DCI>':  Class.DCI}, shapes)
    replace_text({'<DCI_MAJ>':  Class.DCI_MAJ}, shapes)

    replace_text({'<1DR>':  Class.DR1_affichage}, shapes)
    replace_text({'<DPR>':  Class.DPR_affichage}, shapes)
    replace_text({'<DCF>':  Class.DCF_affichage}, shapes)
    replace_text({'<DCF_MAJ>':  Class.DCF_MAJ}, shapes)
    replace_text({'<DEC>':  Class.DEC_affichage}, shapes)
    replace_text({'<F0>':  Class.F0}, shapes)
    replace_text({'<DDCI>':  Class.DDCI_affichage}, shapes)
    replace_text({'<DDCI_MAJ>':  Class.DDCI_MAJ}, shapes)

    replace_text({'<TSJ1>':  Class.TSJ[0]}, shapes)

    # replace_text({'<TSJ2>':  Class.TSJ2}, shapes)
    # replace_text({'<TSJ3>':  Class.TSJ3}, shapes)
    # replace_text({'<TSJ4>':  Class.TSJ4}, shapes)
    # replace_text({'<TSJ5>':  Class.TSJ5}, shapes)
    # replace_text({'<PCS1>':  Class.PCS1}, shapes)
    # replace_text({'<PCS2>':  Class.PCS2}, shapes)
    # replace_text({'<PCS3>':  Class.PCS3}, shapes)
    # replace_text({'<PCS4>':  Class.PCS4}, shapes)
    # replace_text({'<PCS5>':  Class.PCS5}, shapes)
    cpn = (f'{float(Class.CPN):.2f}')
    cpn = str(cpn) + "%"
    cpn = cpn.replace(".", ",")
    replace_text({'<CPN>':  cpn}, shapes)
    pdi = Class.PDI + "%"
    replace_text({'<PDI>':  pdi}, shapes)

    replace_text({'<ADCF>': Class.ADCF_affichage}, shapes) 

    #on rjaoute les % sans que ca ait un impact sur les calculs
    bac = Class.BAC + "%"
    bac = bac.replace(".", ",")

    replace_text({'<BAC>':  bac}, shapes)
    
    replace_text({'<BCPN>':  Class.BCPN}, shapes)

    com = str(Class.COM) + "%"
    com = com.replace(".", ",")
    replace_text({'<COM>':  com}, shapes)
    nsd = Class.NSD + "%"
    replace_text({'<NSD>':  nsd}, shapes)
    nsm = Class.NSM + "%"
    replace_text({'<NSM>':  nsm}, shapes)
    nsf = str(Class.NSF) + "%"
    replace_text({'<NSF>':  nsf}, shapes)

    abdac = str(Class.ABDAC).replace(".", ",")
    replace_text({'<ABDAC>':  abdac}, shapes)
    
    dbac = str(Class.DBAC) + "%"
    dbac = dbac.replace(".", ",")

    replace_text({'<DBAC>':  dbac}, shapes)
    
    deg = str(Class.DEG)
    deg = deg.replace(".", ",")
    replace_text({'<DEG>':  deg}, shapes)

def elementsToReplaceCalcul(Class, shapes):
    replace_text({'<balisedeg1>': Class.balisedeg}, shapes) 
    replace_text({'<balisedeg2>': Class.balisedeg2}, shapes)     
    replace_text({'<balisedeg3>': Class.balisedeg3}, shapes)
    replace_text({'<balisedeg4>': Class.balisedeg4}, shapes) 

    replace_text({'<baliseCM>': Class.baliseCM}, shapes) 
    replace_text({'<baliseCM2>': Class.baliseCM2}, shapes)
    replace_text({'<baliseCM22>': Class.baliseCM22}, shapes) 

    replace_text({'<baliseCM3>': Class.baliseCM3}, shapes) 
    replace_text({'<baliseCM4>': Class.baliseCM4}, shapes)
    replace_text({'<baliseCM5>': Class.baliseCM5}, shapes)
    replace_text({'<baliseCM6>': Class.baliseCM6}, shapes)

    replace_text({'<SV>': Class.SV}, shapes) 

    replace_text({'<DDCI>': Class.DDCI}, shapes) 

    replace_text({'<BLOCDIVIDENDE>': Class.BLOCDIVIDENDE}, shapes)
    replace_text({'<F0s>':  Class.F0s}, shapes)


    replace_text({'<1PDC>': Class.PDC1_affichage}, shapes) 
    
    replace_text({'<2PDC>': Class.PDC2_affichage}, shapes) 
    replace_text({'<DDR>': Class.DDR}, shapes)
    replace_text({'<DIC>': Class.DIC}, shapes)
    
    pdiperf = str(int(Class.PDIPERF)) + "%"
    pdiperf = pdiperf.replace(".", ",")

    replace_text({'<PDIPERF>': pdiperf}, shapes)
    replace_text({'<1PR>': Class.PR1}, shapes)
    replace_text({'<DPRR>': Class.DPRR}, shapes)
    replace_text({'<TDP>': Class.TDP}, shapes)
    replace_text({'<GC>': Class.GC}, shapes)
    gca = str(Class.GCA) + "%"
    gca = gca.replace(".", ",")

    replace_text({'<GCA>': gca}, shapes)
    gce = ("{:.2f}".format(Class.GCE))
    gce = str(gce) + "%"
    gce = gce.replace(".", ",")

    replace_text({'<GCE>': gce}, shapes)
    replace_text({'<ABAC>': Class.ABAC}, shapes)
    replace_text({'<NDR>': Class.NDR }, shapes)
    replace_text({'<ADPR>': Class.ADPR}, shapes)
    replace_text({'<SJR1>': Class.SJR1}, shapes)
    replace_text({'<SJR2>': Class.SJR2}, shapes)
    replace_text({'<SJR3>': Class.SJR3}, shapes)
    replace_text({'<SJR4>': Class.SJR4}, shapes)
    replace_text({'<SJR5>': Class.SJR5}, shapes)
    replace_text({'<SJR6>': Class.SJR6}, shapes)
    replace_text({'<SJR8>': Class.SJR8}, shapes)

    replace_text({'<SJR7>': Class.SJR7}, shapes)

    replace_text({'<F1>': Class.F1}, shapes)
    replace_text({'<F2>': Class.F2}, shapes)

    replace_text({'<TDS>': Class.TDS}, shapes)
    replace_text({'<DU>': Class.DU}, shapes)
    replace_text({'<DU1>':  Class.DU1}, shapes)

    replace_text({'<sponsor>': Class.sponsor}, shapes)
    replace_text({'<SPONSOR>': Class.SPONSOR}, shapes)

    replace_text({'<TICKER>': Class.TICKER}, shapes)
    replace_text({'<NOMSOUSJACENT>': Class.NOMSOUSJACENT}, shapes)
    replace_text({'<DIVIDENDE>': Class.DIVIDENDE}, shapes)
    replace_text({'<SITE>': Class.Site}, shapes)
    replace_text({'<DPCI>': Class.DPCI}, shapes) 
    replace_text({'<DPCI_MAJ>':  Class.DPCI_MAJ}, shapes)

    replace_text({'<NOMP1>': Class.NOMP1}, shapes) 
    replace_text({'<NOMSOUSJACENTP1>': Class.NOMSOUSJACENTP1}, shapes) 
    replace_text({'<SJR6P1>': Class.SJR6P1}, shapes) 
    replace_text({'<type_bar2>': Class.type_bar2}, shapes) 
    replace_text({'<1PR-1>': Class.PR1_1}, shapes) 
    replace_text({'<ABAC2>': Class.ABAC2}, shapes) 
    replace_text({'<DDP>': Class.DDP}, shapes) 
    replace_text({'<Mémoire>': Class.Memoire}, shapes) 
    replace_text({'<Mémoire2>': Class.Memoire2}, shapes) 
    replace_text({'<Mémoire3>': Class.Memoire3}, shapes) 
    replace_text({'<Mémoire4>': Class.Memoire4}, shapes)
    replace_text({'<Mémoire5>': Class.Memoire5}, shapes) 
    replace_text({'<Mémoire6>': Class.Memoire6}, shapes) 
    bfp = str(Class.BFP) + "%"
    replace_text({'<BFP>': bfp}, shapes) 
    replace_text({'<PAGE>': Class.PAGE}, shapes) 

    replace_text({'<DDPP>': Class.DDPP}, shapes) 

    TRA_replace(Class, shapes)

    #CHANGER LE NOM DE LA BALISE ET DE LA CLASS dans myclass.py
    replace_text({'<balise>': Class.balise}, shapes)

    hardcode_replace(Class, shapes)

def hardcode_replace(Class, shapes):
    replace_text({"l' année": "l'année"}, shapes)
    replace_text({" , ": ", "}, shapes)
    replace_text({"(1)": "⁽¹⁾"}, shapes)
    replace_text({"(2)": "⁽²⁾"}, shapes)

def TRA_replace(Class, shapes):
    replace_text({'<TRA.D.A>': Class.TRA_A_S1}, shapes)
    replace_text({'<TRA.M.A>': Class.TRA_A_S2_100}, shapes)
    replace_text({'<TRA.MG.A>': Class.TRA_A_S2_GAIN}, shapes)
    replace_text({'<TRA.M.SJ>': Class.TRA_M_SJ}, shapes)
    replace_text({'<TRA.F.A>': Class.TRA_F_A}, shapes)
    replace_text({'<TRA.F.SJ>': Class.TRA_F_SJ}, shapes)
    replace_text({'<TRA.MRA.MIN.A>': Class.TRA_MIN_A}, shapes)
    replace_text({'<TRA.ECHEANCE.PERTE.A>': Class.TRA_echeance_perte_A}, shapes)
    replace_text({'<BALISECMTRA>': Class.TRA_A_S2_100}, shapes)

    replace_text({'<TRA_D_P>': Class.TRA_D_P}, shapes)
    replace_text({'<TRA_M_P>': Class.TRA_M_P}, shapes)
    replace_text({'<TRA_M_PM>': Class.TRA_M_PM}, shapes)
    replace_text({'<TRA_GM_P>': Class.TRA_GM_P}, shapes)
    replace_text({'<TRA_GM_PM>': Class.TRA_GM_PM}, shapes)
    replace_text({'<TRA_GM_PM>': Class.TRA_F_P}, shapes)

    replace_text({'<TRA_MRA_MIN_P>': Class.TRA_MRA_MIN_P}, shapes)
    replace_text({'<TRA_MRA_MIN_PM>': Class.TRA_MRA_MIN_PM}, shapes)
    replace_text({'<TRA_TOUT-1_P>': Class.TRA_TOUT_1_P}, shapes)
    replace_text({'<TRA_MRA_MAX_P>': Class.TRA_MRA_P}, shapes)
    replace_text({'<TRA_MRD_P>': Class.TRA_MRD_P}, shapes)
    replace_text({'<TRA_TOUT_SAUF_P>': Class.TRA_TOUT_SAUF_P}, shapes)
    replace_text({'<TRA_TOUT_P>': Class.TRA_TOUT_P}, shapes)
    replace_text({'<TRA_MRE_MIN_P>': Class.MRE_MIN_P}, shapes)



    Dates_maj(Class, shapes)

def Dates_maj(Class, shapes):
    replace_text({'<1PDC_MAJ>': Class.PDC1_MAJ}, shapes)
    replace_text({'<2PDC_MAJ>': Class.PDC2_MAJ}, shapes)
    replace_text({'<DDR_MAJ>': Class.DDR_MAJ}, shapes)
    replace_text({'<DEC_MAJ>': Class.DEC_MAJ}, shapes)
    Replace_Boucle_Dates(Class, shapes)

def Replace_Boucle_Dates(Class, shapes):
    replace_text({'<Datesconstatations1>': Class.Datesconstatations1}, shapes)
    replace_text({'<Datesconstatations2>': Class.Datesconstatations2}, shapes)
    replace_text({'<Datesconstatations3>': Class.Datesconstatations3}, shapes)
    replace_text({'<Datesconstatations4>': Class.Datesconstatations4}, shapes)

    replace_text({'<Datesremb1>': Class.Datesremb1}, shapes)
    replace_text({'<Datesremb2>': Class.Datesremb2}, shapes)
    replace_text({'<Datesremb3>': Class.Datesremb3}, shapes)
    replace_text({'<Datesremb4>': Class.Datesremb4}, shapes)
    replace_text({'<Datesremb5>': Class.Datesremb5}, shapes)
    replace_text({'<Datesremb6>': Class.Datesremb6}, shapes)
    replace_text({'<Datesremb7>': Class.Datesremb7}, shapes)
    replace_text({'<Datesremb8>': Class.Datesremb8}, shapes)

    replace_text({'<Datespaiement1>': Class.Datespaiement1}, shapes)
    replace_text({'<Datespaiement2>': Class.Datespaiement2}, shapes)
    replace_text({'<Datespaiement3>': Class.Datespaiement3}, shapes)
    replace_text({'<Datespaiement4>': Class.Datespaiement4}, shapes)
    replace_text({'<Datespaiement5>': Class.Datespaiement5}, shapes)
    replace_text({'<Datespaiement6>': Class.Datespaiement6}, shapes)
    replace_text({'<Datespaiement7>': Class.Datespaiement7}, shapes)
    replace_text({'<Datespaiement8>': Class.Datespaiement8}, shapes)

#PREMIER RAPPEL A DATE DERNIER RAPPEL

#avoir toutes les slides
def getAllSlides(prs):
    slides = [slide for slide in prs.slides]
    shapes = []

    for slide in slides:
        for shape in slide.shapes:
            shapes.append(shape)
    return shapes



def add_row(table) -> _Row:
    new_row = deepcopy(table._tbl.tr_lst[-1]) 
    # duplicating last row of the table as a new row to be added

    for tc in new_row.tc_lst:
        cell = _Cell(tc, new_row.tc_lst)
        cell.text = '' # defaulting cell contents to empty text

    table._tbl.append(new_row) 
    return table

    
#remplace le texte sur le ppt 
def replace_text(replacements: dict, shapes: List):
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:

                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = cur_text.replace(str(match), str(replacement))
                            run.text = new_text


            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if match in cell.text:
                            new_text = cell.text.replace(match, replacement)
                            cell.text = new_text


#load le ppt, remplace les balises et le sauvegarde
def ChangeTextOnPpt(Class):
    NAME = "result/"+ Class.Nom + "- " + Class.Isin + "result.pptx" 
    prs_string = "templates/"+ Class.template + ".pptx" 
    prs = Presentation(prs_string)
    shapes = getAllSlides(prs)
    Class.shapes = shapes

    elementsToReplaceDegressivite(Class, shapes)

    elementsToReplaceRemplacement(Class, shapes)
    elementsToReplaceCalcul(Class, shapes)

    #on repasse une 2 eme fois au cas ou certaines transformations soient mal placées
    elementsToReplaceDegressivite(Class, shapes)

    elementsToReplaceRemplacement(Class, shapes)
    elementsToReplaceCalcul(Class, shapes)

    compteur = 0
    for shape in Class.shapes:
        textbox = shape
        for i in Class.deleteblocs:

            try:
                if (i in textbox.text):
                        print("paragraphe supprimé")
                        compteur+=1
                        sp = textbox.element
                        sp.getparent().remove(sp)

            except Exception:
                pass
            
    if compteur > 1:
        print(compteur, "paragraphes ont été supprimés")
    else:

        print(compteur, "paragraphe a été supprimé")
    # from colormath.color_objects import sRGBColor
    compteur = 0
    compteur_tab = 1
    
    tbl = None


    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                numberoftickers = len(Class.Yahoo) -1 #pour savoir le nombre de rows a ajouter
                if(compteur_tab == 1 and len(Class.Yahoo) > 1):  #On ajoute des rows pour le tableau de performance des tickers yahoo
                    for newcolumn in range((numberoftickers)):
                        table = add_row(table)

                for row_idx, row in enumerate(table.rows): 
                    for col_idx, cell in enumerate(row.cells):

                        if (compteur_tab == 1 and col_idx == 0 and row_idx > 0): #on ajoute les sousjacent a la premiere colonne de chaque ligne
                            cell = table.cell(row_idx, col_idx)
                            cell.text = str(Class.Yahoo_value_name[row_idx -1]) +" " + str(Class.Yahoo_value_dividende[row_idx -1])


                        # print("%r is cells[%d][%d]" % (cell, row_idx, col_idx))
                        cell = table.cell(row_idx, col_idx)
                        # cell.text = 'TEST DE LA STREET'
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(8)
                        # paragraph.font.color = 'rgb(0, 0, 0)'

                        # paragraph.font.color.rgb = screen.bgcolor(color)

                        if(compteur_tab ==1 and row_idx >= 1 and col_idx >= 1 ):
                            cell.text = Class.Yahoo_value[row_idx - 1][col_idx - 1]
                            paragraph = cell.text_frame.paragraphs[0]
                            paragraph.font.size = Pt(8)
                        
                compteur_tab +=1



            if shape.has_text_frame:
                if ("<graph1>" in shape.text):
                    cur_text = shape.text
                    new_text = cur_text.replace(str("<graph1>"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph1.png", Inches(0), Inches(6.75), Inches(7))
                
                if ("<graph2>" in shape.text):  
                    cur_text = shape.text
                    new_text = cur_text.replace(str("<graph2>"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph_scenario_def.png", Inches(0.45), Inches(1.75), Inches(3.6))

                if ("<graph3>" in shape.text):  
                    cur_text = shape.text
                    new_text = cur_text.replace(str("<graph3>"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph_scenario_median.png", Inches(0.45), Inches(4.78), Inches(3.6))
    
                if ("<graph4>" in shape.text):  
                    cur_text = shape.text
                    new_text = cur_text.replace(str("<graph4>"), str(""))
                    shape.text = new_text
                    pic = slide.shapes.add_picture("graph_scenario_fav.png", Inches(0.45), Inches(7.68), Inches(3.45))
                   
                if ("<graph5>" in shape.text):  
                        cur_text = shape.text
                        new_text = cur_text.replace(str("<graph5>"), str(""))
                        shape.text = new_text
                        pic = slide.shapes.add_picture("graph5.png", Inches(0), Inches(4.8), width=Inches(6.8), height=Inches(4))


    try:
        print("Nettoyage du projet, supression des documents inutiles")
        os.remove("graph1.png")
        print("Suppression de graph1.png")
        os.remove("graph2.png")
        print("Suppression de graph2.png")
        os.remove("graph3.png")
        print("Suppression de graph3.png")
        os.remove("graph4.png")
        print("Suppression de graph4.png")
        os.remove("graph5.png")
        print("Suppression de graph5.png")

    except Exception:
        print("OH NONNNN")
    prs.save(NAME)






#balise.py     if (strike == "best strike"): rajouter à la fin si cest de l'action, ou de l'indice ou des actions ou des indices   + NOM
