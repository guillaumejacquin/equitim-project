from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import plotly.express as px 

# create presentation with 1 slide ------
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

# define chart data ---------------------

chart_data = CategoryChartData()
chart_data.categories = ['']
chart_data.categories = ['t1', 't5', 't12']

chart_data.add_series('West',(100, 80, 60.7))

x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
).chart

chart.has_legend = True
chart.legend.include_in_layout = False
chart.series[0].smooth = True
chart.value_axis.minimum_scale = 0



prs.save('chart-02.pptx')




  
  
# Creating the Figure instance
  
# showing the plot
